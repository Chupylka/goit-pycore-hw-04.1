from pptx import Presentation 
from pptx.util import Inches, Pt 
from pptx.dml.color import RGBColor 
 
# Створення презентації 
prs = Presentation() 
title_slide_layout = prs.slide_layouts[0] 
content_slide_layout = prs.slide_layouts[1] 
 
# Функція для створення слайдів 
def add_slide(title, content): 
    slide = prs.slides.add_slide(content_slide_layout) 
    slide.shapes.title.text = title 
    content_box = slide.placeholders[1] 
    content_box.text = content 
 
# Слайд 1: Вступ 
slide = prs.slides.add_slide(title_slide_layout) 
slide.shapes.title.text = "Прогнозування ціни оренди житла" 
slide.placeholders[1].text = "Презентація фінального проєкту\n[Твоє ім’я]" 
 
# Слайд 2–3: Аналіз та підготовка даних 
add_slide("Аналіз даних",  
          "- Всього записів: 4746\n" 
          "- Жодного пропущеного значення\n" 
          "- Найвища кореляція: Size ↔ Rent (0.6)\n" 
          "- Виявлено викиди (до 3.5 млн)") 
 
add_slide("Підготовка даних",  
          "- Перетворення 'Floor' на числовий рівень\n" 
          "- One-hot кодування: Area Type, City тощо\n" 
          "- Видалено текстові колонки\n" 
          "- Всього ознак після обробки: понад 20") 
 
# Слайд 4: Моделювання 
add_slide("Моделювання",  
          "- Модель 1: Лінійна регресія\n" 
          "- Модель 2: Random Forest\n" 
          "- Розділення: 80% тренування / 20% тест\n" 
          "- Метрики: MAE, RMSE, R² Score") 
 
# Слайд 5: Результати - Лінійна регресія 
add_slide("Лінійна регресія",  
          "- MAE: ~10,000\n" 
          "- RMSE: ~15,000–20,000\n" 
          "- R² Score: ~0.60\n" 
          "- Недооцінює дорогі квартири") 
 
# Слайд 6: Результати - Random Forest 
add_slide("Random Forest",  
          "- MAE: ~4,000–6,000\n" 
          "- RMSE: ~7,000–9,000\n" 
          "- R² Score: ~0.85\n" 
          "- Добре працює з нелінійними залежностями") 
 
# Слайд 7: Висновки та інсайти 
add_slide("Інсайти та враження",  
          "- Якість вхідних ознак критично важлива\n" 
          "- Складні моделі точніші, але потребують ресурсів\n" 
          "- Цей проєкт покращив мої навички аналізу та моделювання") 
 
# Слайд 8: Подяка 
add_slide("Дякую за увагу!", "Готовий відповісти на ваші запитання 🙂") 
 
# Збереження презентації 
pptx_path = "House_Rent_Prediction_Presentation.pptx"
prs.save(pptx_path) 
 
pptx_path

print("Презентацію збережено")

import os
os.system("open House_Rent_Prediction_Presentation.pptx")