from pptx import Presentation 
from pptx.util import Inches 
 
# Створення презентації 
prs = Presentation() 
 
# Додавання титульного слайда 
slide_layout = prs.slide_layouts[0]  # Титульний слайд 
slide = prs.slides.add_slide(slide_layout) 
title = slide.shapes.title 
subtitle = slide.placeholders[1] 
title.text = "Прогнозування ціни оренди нерухомості" 
subtitle.text = "Фінальний проєкт з Mathematics for Computer Science" 
 
# Вступний слайд 
slide_layout = prs.slide_layouts[1]  # Заголовок + текст 
slide = prs.slides.add_slide(slide_layout) 
title = slide.shapes.title 
content = slide.placeholders[1] 
title.text = "Вступ" 
content.text = "Я - дослідник даних, що аналізує ціни оренди житла та прогнозує їх на основі різних факторів." 
 
# Аналіз та підготовка даних - 1 
slide = prs.slides.add_slide(slide_layout) 
title = slide.shapes.title 
content = slide.placeholders[1] 
title.text = "Аналіз даних" 
content.text = "• Оцінка діапазонів значень\n• Виявлення кореляцій між параметрами\n• Перевірка на пропущені значення" 
 
# Аналіз та підготовка даних - 2 
slide = prs.slides.add_slide(slide_layout) 
title = slide.shapes.title 
content = slide.placeholders[1] 
title.text = "Підготовка даних" 
content.text = "• Заповнення пропущених значень\n• Нормалізація числових ознак\n• Закодування категоріальних змінних" 
 
# Моделювання 
slide = prs.slides.add_slide(slide_layout) 
title = slide.shapes.title 
content = slide.placeholders[1] 
title.text = "Моделювання" 
content.text = "• Використані 2 моделі регресії\n• Оцінка точності моделей за метриками RMSE та MAE\n• Використання крос-валідації" 
 
# Аналіз результатів - 1 
slide = prs.slides.add_slide(slide_layout) 
title = slide.shapes.title 
content = slide.placeholders[1] 
title.text = "Аналіз результатів" 
content.text = "• Побудовано графіки порівняння справжніх і передбачених значень\n• Аналіз помилок прогнозування" 
 
# Аналіз результатів - 2 
slide = prs.slides.add_slide(slide_layout) 
title = slide.shapes.title 
content = slide.placeholders[1] 
title.text = "Порівняння моделей" 
content.text = "• Визначено найкращу модель за точністю\n• Аналіз прикладів із найбільшими помилками" 
 
# Висновки та рефлексія 
slide = prs.slides.add_slide(slide_layout) 
title = slide.shapes.title 
content = slide.placeholders[1] 
title.text = "Висновки та рефлексія" 
content.text = "• Отримані цінні інсайти щодо факторів, що впливають на оренду\n• Використання різних моделей дало змогу оцінити їх ефективність\n• Нові навички у роботі з даними та побудові прогнозних моделей" 
 
# Додавання завершального слайда 
slide = prs.slides.add_slide(prs.slide_layouts[5])  # Титульний слайд без підзаголовка 
title = slide.shapes.title 
title.text = "Дякую за увагу!" 
 
# Збереження презентації 
pptx_path = "/Users/vitalinka/vscode-basics/data/Final_Project_Presentation.pptx"
prs.save(pptx_path) 

print("Starting script")

from pptx import Presentation

prs = Presentation()

# Додайте інші повідомлення для перевірки
print("Presentation object created")

pptx_path = "/Users/vitalinka/vscode-basics/data/Final_Project_Presentation.pptx"
prs.save(pptx_path)

print(f"Presentation saved at {pptx_path}")

pptx_path